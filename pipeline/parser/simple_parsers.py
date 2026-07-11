"""Word / HTML / PPT / 纯文本 解析器。"""

import logging

from pipeline.parser.base import BaseParser, ParseResult, RawDocument

logger = logging.getLogger(__name__)


class WordParser(BaseParser):
    MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    def supports(self, mime_type: str) -> bool:
        return mime_type == self.MIME

    def parse(self, doc: RawDocument) -> ParseResult:
        try:
            from docx import Document

            d = Document(doc.file_path)
            paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
            tables_html = [self._table_to_html(t) for t in d.tables]
            return ParseResult(
                success=True,
                markdown="\n\n".join(paragraphs),
                tables=tables_html,
                metadata={"parser": "WordParser"},
                confidence=0.9,
            )
        except Exception as e:  # noqa: BLE001
            return ParseResult(success=False, markdown="", error=str(e))

    @staticmethod
    def _table_to_html(table) -> str:
        rows_html = []
        for row in table.rows:
            cells = "".join(f"<td>{c.text}</td>" for c in row.cells)
            rows_html.append(f"<tr>{cells}</tr>")
        return f"<table>{''.join(rows_html)}</table>"


class HTMLParser(BaseParser):
    def supports(self, mime_type: str) -> bool:
        return mime_type in ("text/html", "application/xhtml+xml")

    def parse(self, doc: RawDocument) -> ParseResult:
        try:
            from bs4 import BeautifulSoup

            html = open(doc.file_path, encoding="utf-8", errors="ignore").read()
            soup = BeautifulSoup(html, "lxml")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            text = soup.get_text("\n", strip=True)
            tables = [str(t) for t in soup.find_all("table")]
            return ParseResult(
                success=True, markdown=text, tables=tables,
                metadata={"parser": "HTMLParser"}, confidence=0.85,
            )
        except Exception as e:  # noqa: BLE001
            return ParseResult(success=False, markdown="", error=str(e))


class PPTParser(BaseParser):
    MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def supports(self, mime_type: str) -> bool:
        return mime_type == self.MIME

    def parse(self, doc: RawDocument) -> ParseResult:
        try:
            from pptx import Presentation

            prs = Presentation(doc.file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if shape.has_text_frame and shape.text_frame.text.strip()
                ]
                if texts:
                    slides_text.append(f"## 第{i}页\n" + "\n".join(texts))
            return ParseResult(
                success=True, markdown="\n\n".join(slides_text),
                metadata={"parser": "PPTParser", "slide_count": len(prs.slides)},
                confidence=0.9,
            )
        except Exception as e:  # noqa: BLE001
            return ParseResult(success=False, markdown="", error=str(e))


class PlainTextParser(BaseParser):
    """OCR 转写文本 / txt 直接入库"""

    def supports(self, mime_type: str) -> bool:
        return mime_type == "text/plain"

    def parse(self, doc: RawDocument) -> ParseResult:
        try:
            text = open(doc.file_path, encoding="utf-8", errors="ignore").read()
            return ParseResult(
                success=True, markdown=text,
                metadata={"parser": "PlainTextParser"}, confidence=1.0,
            )
        except Exception as e:  # noqa: BLE001
            return ParseResult(success=False, markdown="", error=str(e))
